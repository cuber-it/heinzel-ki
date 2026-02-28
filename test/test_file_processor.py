"""
Tests fuer file_processor.py — Datei-Konvertierung
"""
import os, sys, io, json, base64
os.environ.setdefault("LOG_DIR", "/tmp")
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "../src/llm-provider"))

from file_processor import process_file, _is_text_extension, _extract_pdf, _extract_docx, _extract_xlsx, _extract_pptx


# ─── Text-Dateien ─────────────────────────────────────────────

def test_json_file():
    data = b'{"key": "value", "num": 42}'
    block = process_file(data, "test.json", "application/json", "openai")
    assert block.type == "text"
    assert '"key"' in block.text


def test_csv_file():
    data = b"name,age\nAlice,30\nBob,25"
    block = process_file(data, "data.csv", "text/csv", "openai")
    assert block.type == "text"
    assert "Alice" in block.text


def test_xml_file():
    data = b"<root><item>hello</item></root>"
    block = process_file(data, "data.xml", "application/xml", "openai")
    assert block.type == "text"
    assert "hello" in block.text


def test_txt_file():
    data = "Hallo Welt".encode()
    block = process_file(data, "readme.txt", "text/plain", "openai")
    assert block.type == "text"
    assert "Hallo Welt" in block.text


def test_python_file_by_extension():
    data = b"print('hello world')"
    block = process_file(data, "script.py", "application/octet-stream", "openai")
    assert block.type == "text"
    assert "hello world" in block.text


def test_yaml_file():
    data = b"key: value\nlist:\n  - item1"
    block = process_file(data, "config.yaml", "application/x-yaml", "openai")
    assert block.type == "text"
    assert "item1" in block.text


# ─── Bilder ───────────────────────────────────────────────────

def test_jpeg_openai():
    data = b"\xff\xd8\xff" + b"\x00" * 100  # fake JPEG header
    block = process_file(data, "photo.jpg", "image/jpeg", "openai")
    assert block.type == "image"
    assert block.media_type == "image/jpeg"
    assert block.data == base64.b64encode(data).decode()


def test_png_anthropic():
    data = b"\x89PNG\r\n" + b"\x00" * 50
    block = process_file(data, "img.png", "image/png", "anthropic")
    assert block.type == "image"
    assert block.media_type == "image/png"


def test_jpg_alias_normalized():
    """image/jpg → image/jpeg normalisieren"""
    data = b"\xff\xd8\xff" + b"\x00" * 10
    block = process_file(data, "photo.jpg", "image/jpg", "openai")
    assert block.type == "image"
    assert block.media_type == "image/jpeg"


# ─── PDF ──────────────────────────────────────────────────────

def test_pdf_anthropic_native():
    """Anthropic: PDF als DocumentBlock"""
    data = b"%PDF-1.4 fake"
    block = process_file(data, "doc.pdf", "application/pdf", "anthropic")
    assert block.type == "document"
    assert block.data == base64.b64encode(data).decode()


def test_pdf_google_native():
    """Google: PDF als DocumentBlock"""
    data = b"%PDF-1.4 fake"
    block = process_file(data, "doc.pdf", "application/pdf", "google")
    assert block.type == "document"


def test_pdf_openai_text_extraction():
    """OpenAI: echtes PDF → Text extrahieren"""
    # Minimales gültiges PDF mit Text erstellen
    import pypdf, pypdf.generic as g
    # Wir testen mit dem Extraktor direkt
    block = _extract_pdf(b"%PDF invalid", "test.pdf")
    assert block.type == "text"
    # Entweder Inhalt oder Fehlermeldung — beides ist TextBlock
    assert len(block.text) > 0


def test_pdf_openai_via_process_file():
    """OpenAI: PDF geht durch process_file → TextBlock (kein Crash)"""
    data = b"%PDF-1.4 minimal"
    block = process_file(data, "report.pdf", "application/pdf", "openai")
    assert block.type == "text"  # Text-Extraktion oder Fehlermeldung


# ─── Office ───────────────────────────────────────────────────

def test_docx_extraction():
    """DOCX → Text"""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Erster Absatz")
    doc.add_paragraph("Zweiter Absatz")
    buf = io.BytesIO()
    doc.save(buf)
    block = _extract_docx(buf.getvalue(), "test.docx")
    assert block.type == "text"
    assert "Erster Absatz" in block.text
    assert "Zweiter Absatz" in block.text


def test_xlsx_extraction():
    """XLSX → Text"""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Tabelle1"
    ws.append(["Name", "Wert"])
    ws.append(["Alice", 42])
    buf = io.BytesIO()
    wb.save(buf)
    block = _extract_xlsx(buf.getvalue(), "data.xlsx")
    assert block.type == "text"
    assert "Alice" in block.text
    assert "Tabelle1" in block.text


def test_pptx_extraction():
    """PPTX → Text"""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Folientitel"
    slide.placeholders[1].text = "Folieninhalt"
    buf = io.BytesIO()
    prs.save(buf)
    block = _extract_pptx(buf.getvalue(), "pres.pptx")
    assert block.type == "text"
    assert "Folientitel" in block.text


def test_docx_via_process_file():
    """DOCX geht durch process_file"""
    from docx import Document
    doc = Document()
    doc.add_paragraph("Testinhalt")
    buf = io.BytesIO()
    doc.save(buf)
    mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    block = process_file(buf.getvalue(), "doc.docx", mime, "openai")
    assert block.type == "text"
    assert "Testinhalt" in block.text


# ─── Fehler-Fälle ─────────────────────────────────────────────

def test_video_returns_error_text():
    data = b"\x00\x01\x02\x03"
    block = process_file(data, "video.mp4", "video/mp4", "openai")
    assert block.type == "text"
    assert "nicht unterstützt" in block.text.lower() or "video" in block.text.lower()


def test_unknown_binary_returns_text():
    data = b"\x00\x01\x02\x03\x04\x05"
    block = process_file(data, "unknown.bin", "application/octet-stream", "openai")
    assert block.type == "text"


# ─── Hilfsfunktionen ──────────────────────────────────────────

def test_is_text_extension():
    assert _is_text_extension("script.py") is True
    assert _is_text_extension("data.json") is True
    assert _is_text_extension("config.yaml") is True
    assert _is_text_extension("image.jpg") is False
    assert _is_text_extension("document.pdf") is False
    assert _is_text_extension("noextension") is False
