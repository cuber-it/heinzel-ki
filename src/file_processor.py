"""
H.E.I.N.Z.E.L. Provider — File Processor

Konvertiert hochgeladene Dateien in ContentBlocks.
Jeder Provider deklariert seine nativen MIME-Types —
alles andere wird soweit möglich in Text oder Bilder konvertiert.

Unterstuetzte Strategien:
  native    → direkt als ImageBlock / DocumentBlock
  text      → Inhalt als TextBlock (JSON, XML, CSV, Code, ...)
  extract   → Text-Extraktion (PDF, DOCX, XLSX, PPTX)
  error     → klare Fehlermeldung als TextBlock (Video, Audio, ...)
"""
from __future__ import annotations

import base64
import io
from dataclasses import dataclass, field
from typing import Literal

from models import TextBlock, ImageBlock, DocumentBlock, ContentBlock


# ─── Provider-Fähigkeiten ─────────────────────────────────────────────────────

# MIME-Types die ein Provider nativ als Binary-Block verarbeiten kann
NATIVE_IMAGE_TYPES = {"image/jpeg", "image/jpg", "image/png", "image/gif", "image/webp"}
NATIVE_PDF_TYPES   = {"application/pdf"}

PROVIDER_NATIVE: dict[str, set[str]] = {
    "anthropic": NATIVE_IMAGE_TYPES | NATIVE_PDF_TYPES,
    "google":    NATIVE_IMAGE_TYPES | NATIVE_PDF_TYPES,
    "openai":    NATIVE_IMAGE_TYPES,  # OpenAI: kein natives PDF
}

# Text-basierte MIME-Types → immer als TextBlock
TEXT_MIME_TYPES = {
    "text/plain", "text/html", "text/markdown", "text/csv",
    "text/xml", "application/xml",
    "application/json", "application/javascript",
    "application/x-yaml", "text/yaml",
    # Code
    "text/x-python", "text/x-java-source", "text/x-c", "text/x-c++",
    "text/x-shellscript", "application/x-sh",
    "text/x-sql",
}

# Office-Formate → Text-Extraktion
OFFICE_MIME_TYPES = {
    "application/pdf",  # Fallback wenn Provider kein natives PDF kann
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",   # docx
    "application/msword",                                                          # doc (best-effort)
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",          # xlsx
    "application/vnd.ms-excel",                                                    # xls
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",  # pptx
    "application/vnd.ms-powerpoint",                                               # ppt
}

# Nicht unterstützte Typen → Error-TextBlock
UNSUPPORTED_MIME_TYPES = {
    "video/", "audio/",
    "application/octet-stream",
    "application/x-executable",
}


# ─── Haupt-Konverter ──────────────────────────────────────────────────────────

def process_file(
    data: bytes,
    filename: str,
    mime_type: str,
    provider_name: str,
) -> ContentBlock:
    """
    Konvertiert eine Datei in einen ContentBlock passend für den Provider.

    Priorität:
    1. Natives Binary (Bild/PDF wenn Provider unterstützt)
    2. Text-basierte MIME-Types direkt als TextBlock
    3. Office/PDF → Text-Extraktion
    4. Fallback: Fehlermeldung
    """
    mime = (mime_type or "application/octet-stream").lower().split(";")[0].strip()
    native_types = PROVIDER_NATIVE.get(provider_name, NATIVE_IMAGE_TYPES)

    # 1. Nativ: Bilder
    if mime in NATIVE_IMAGE_TYPES and mime in native_types:
        return ImageBlock(
            media_type=mime if mime != "image/jpg" else "image/jpeg",
            data=base64.b64encode(data).decode(),
        )

    # 2. Nativ: PDF (nur wenn Provider es kann)
    if mime == "application/pdf" and mime in native_types:
        return DocumentBlock(data=base64.b64encode(data).decode())

    # 3. Text-basierte Dateien → direkt als Text
    if mime in TEXT_MIME_TYPES or _is_text_extension(filename):
        try:
            text = data.decode("utf-8", errors="replace")
            label = f"[{filename}]\n"
            return TextBlock(text=label + text)
        except Exception as e:
            return TextBlock(text=f"[Fehler beim Lesen von {filename}: {e}]")

    # 4. Office / PDF → Text-Extraktion
    if mime == "application/pdf":
        return _extract_pdf(data, filename)

    if mime in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    ):
        return _extract_docx(data, filename)

    if mime in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
    ):
        return _extract_xlsx(data, filename)

    if mime in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    ):
        return _extract_pptx(data, filename)

    # 5. Nicht unterstützt
    for prefix in UNSUPPORTED_MIME_TYPES:
        if mime.startswith(prefix):
            return TextBlock(
                text=f"[{filename}] Dieser Dateityp ({mime}) wird von keinem Provider unterstützt."
            )

    # 6. Unbekannter Typ → best-effort Text-Versuch, sonst Error
    try:
        text = data.decode("utf-8", errors="strict")
        return TextBlock(text=f"[{filename}]\n{text}")
    except UnicodeDecodeError:
        return TextBlock(
            text=(
                f"[{filename}] Unbekannter Dateityp ({mime}). "
                f"Dateigröße: {len(data)} Bytes. "
                f"Dieser Typ kann nicht verarbeitet werden."
            )
        )


# ─── Extraktoren ──────────────────────────────────────────────────────────────

def _extract_pdf(data: bytes, filename: str) -> TextBlock:
    """PDF → Text via pypdf."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"--- Seite {i} ---\n{text.strip()}")
        if not pages:
            return TextBlock(
                text=f"[{filename}] PDF konnte nicht als Text extrahiert werden "
                     f"(möglicherweise rein bildbasiert). Bitte einen Provider mit "
                     f"nativem PDF-Support verwenden (Anthropic, Google)."
            )
        return TextBlock(text=f"[{filename} — PDF-Inhalt]\n\n" + "\n\n".join(pages))
    except ImportError:
        return TextBlock(text=f"[{filename}] pypdf nicht installiert — PDF-Extraktion nicht verfügbar.")
    except Exception as e:
        return TextBlock(text=f"[{filename}] PDF-Extraktion fehlgeschlagen: {e}")


def _extract_docx(data: bytes, filename: str) -> TextBlock:
    """DOCX → Text via python-docx."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n".join(paragraphs)
        return TextBlock(text=f"[{filename} — Word-Dokument]\n\n{text}")
    except ImportError:
        return TextBlock(text=f"[{filename}] python-docx nicht installiert.")
    except Exception as e:
        return TextBlock(text=f"[{filename}] DOCX-Extraktion fehlgeschlagen: {e}")


def _extract_xlsx(data: bytes, filename: str) -> TextBlock:
    """XLSX → CSV-ähnlicher Text via openpyxl."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(c) if c is not None else "" for c in row]
                if any(v.strip() for v in row_vals):
                    rows.append("\t".join(row_vals))
            if rows:
                sheets.append(f"=== Tabelle: {sheet_name} ===\n" + "\n".join(rows))
        text = "\n\n".join(sheets) if sheets else "(leer)"
        return TextBlock(text=f"[{filename} — Excel]\n\n{text}")
    except ImportError:
        return TextBlock(text=f"[{filename}] openpyxl nicht installiert.")
    except Exception as e:
        return TextBlock(text=f"[{filename}] XLSX-Extraktion fehlgeschlagen: {e}")


def _extract_pptx(data: bytes, filename: str) -> TextBlock:
    """PPTX → Text pro Folie via python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"--- Folie {i} ---\n" + "\n".join(texts))
        text = "\n\n".join(slides) if slides else "(keine Texte gefunden)"
        return TextBlock(text=f"[{filename} — PowerPoint]\n\n{text}")
    except ImportError:
        return TextBlock(text=f"[{filename}] python-pptx nicht installiert.")
    except Exception as e:
        return TextBlock(text=f"[{filename}] PPTX-Extraktion fehlgeschlagen: {e}")


# ─── Hilfsfunktionen ──────────────────────────────────────────────────────────

def _is_text_extension(filename: str) -> bool:
    """Erkennt Text-Dateien anhand der Extension wenn MIME-Type unbekannt."""
    text_extensions = {
        ".txt", ".md", ".markdown", ".rst",
        ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf",
        ".xml", ".html", ".htm", ".svg",
        ".csv", ".tsv",
        ".py", ".js", ".ts", ".jsx", ".tsx", ".vue",
        ".java", ".c", ".cpp", ".h", ".cs", ".go", ".rs", ".rb", ".php",
        ".sh", ".bash", ".zsh", ".fish",
        ".sql", ".graphql",
        ".log", ".env",
    }
    ext = "." + filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    return ext in text_extensions
